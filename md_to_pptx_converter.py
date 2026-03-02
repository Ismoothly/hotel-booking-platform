#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Markdown to PowerPoint Converter using python-pptx
Generates fully editable PPTX from Markdown
"""

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class MarkdownToPPT:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        self.slides = []
        
    def parse_markdown(self, content):
        """Parse markdown content into slides"""
        # Remove front matter (---)
        lines = content.split('\n')
        
        # Skip front matter
        start_idx = 0
        if lines[0].strip() == '---':
            for i in range(1, len(lines)):
                if lines[i].strip() == '---':
                    start_idx = i + 1
                    break
        
        # Process content
        current_slide = None
        i = start_idx
        
        while i < len(lines):
            line = lines[i]
            
            # Skip HTML comments
            if '<!--' in line:
                i += 1
                while i < len(lines) and '-->' not in lines[i]:
                    i += 1
                i += 1
                continue
            
            # Check for slide separator (---)
            if line.strip() == '---':
                if current_slide:
                    self.slides.append(current_slide)
                current_slide = {'title': '', 'subtitle': '', 'content': []}
                i += 1
                continue
            
            # Check for headings
            if line.startswith('# ') and not current_slide:
                # Main title (create first slide)
                current_slide = {
                    'title': line.lstrip('# ').strip(),
                    'subtitle': '',
                    'content': []
                }
            elif line.startswith('## '):
                # Secondary heading = new slide
                if current_slide and current_slide['content']:
                    self.slides.append(current_slide)
                current_slide = {
                    'title': line.lstrip('## ').strip(),
                    'subtitle': '',
                    'content': []
                }
            elif line.startswith('### '):
                # Tertiary heading = subsection title
                if current_slide:
                    current_slide['content'].append({
                        'type': 'heading3',
                        'text': line.lstrip('### ').strip()
                    })
            elif line.startswith('**') and line.endswith('**'):
                # Bold text = section header
                if current_slide:
                    current_slide['content'].append({
                        'type': 'bold_text',
                        'text': line.strip('**').strip()
                    })
            elif line.startswith('- ') or line.startswith('* '):
                # Bullet point
                if current_slide:
                    indent = len(line) - len(line.lstrip())
                    level = indent // 2
                    text = line.lstrip('- *').strip()
                    current_slide['content'].append({
                        'type': 'bullet',
                        'text': text,
                        'level': level
                    })
            elif line.startswith('```'):
                # Code block
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                if current_slide and code_lines:
                    current_slide['content'].append({
                        'type': 'code',
                        'text': '\n'.join(code_lines).strip()
                    })
                i += 1
                continue
            elif line.startswith('|'):
                # Table
                table_lines = [line]
                i += 1
                while i < len(lines) and lines[i].startswith('|'):
                    table_lines.append(lines[i])
                    i += 1
                if current_slide and len(table_lines) >= 3:
                    rows = []
                    for tline in table_lines:
                        cells = [cell.strip() for cell in tline.split('|')[1:-1]]
                        if cells and not all('-' in c for c in cells):
                            rows.append(cells)
                    if rows:
                        current_slide['content'].append({
                            'type': 'table',
                            'rows': rows
                        })
                continue
            elif line.strip():
                # Regular text
                if current_slide:
                    text = line.strip()
                    # Skip empty lines and certain patterns
                    if text and not text.startswith('---'):
                        current_slide['content'].append({
                            'type': 'text',
                            'text': text
                        })
            
            i += 1
        
        # Add last slide
        if current_slide:
            self.slides.append(current_slide)
    
    def add_title_slide(self, title, subtitle):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(25, 118, 210)  # Blue
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.word_wrap = True
            
            for line in subtitle.split('\n'):
                if subtitle_frame.paragraphs:
                    p = subtitle_frame.add_paragraph()
                else:
                    p = subtitle_frame.paragraphs[0]
                p.text = line
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
    
    def add_content_slide(self, title, content_items):
        """Add content slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(25, 118, 210)
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        first = True
        for item in content_items:
            if item['type'] == 'bullet':
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item['text']
                p.level = item.get('level', 0)
                p.font.size = Pt(16 - item.get('level', 0) * 2)
                p.space_before = Pt(6)
                p.space_after = Pt(6)
                
            elif item['type'] == 'text':
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item['text']
                p.font.size = Pt(14)
                p.space_before = Pt(3)
                p.space_after = Pt(3)
                
            elif item['type'] == 'heading3':
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(25, 118, 210)
                p.space_before = Pt(8)
                p.space_after = Pt(6)
                
            elif item['type'] == 'bold_text':
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item['text']
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_before = Pt(6)
                p.space_after = Pt(4)
                
            elif item['type'] == 'code':
                if first:
                    p = text_frame.paragraphs[0]
                    first = False
                else:
                    p = text_frame.add_paragraph()
                
                p.text = item['text']
                p.font.name = 'Courier New'
                p.font.size = Pt(10)
                p.font.color.rgb = RGBColor(100, 100, 100)
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                
            elif item['type'] == 'table':
                # Add simple table representation
                if first:
                    first = False
                
                rows = item['rows']
                p = text_frame.add_paragraph() if not first else text_frame.paragraphs[0]
                
                for row in rows:
                    p.text = ' | '.join(row)
                    p.font.size = Pt(12)
                    p.space_after = Pt(3)
                    if row == rows[0]:
                        p.font.bold = True
                    p = text_frame.add_paragraph()
    
    def generate(self, markdown_file, output_file):
        """Generate PPTX from Markdown file"""
        with open(markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()
        
        self.parse_markdown(content)
        
        # Add slides
        for i, slide_data in enumerate(self.slides):
            if i == 0:
                # First slide is title slide
                subtitle_lines = []
                for item in slide_data['content']:
                    if item['type'] == 'bullet':
                        subtitle_lines.append('• ' + item['text'])
                    elif item['type'] == 'text':
                        subtitle_lines.append(item['text'])
                
                self.add_title_slide(slide_data['title'], '\n'.join(subtitle_lines[:3]))
            else:
                # Content slides
                self.add_content_slide(slide_data['title'], slide_data['content'])
        
        # Save
        self.prs.save(output_file)
        print(f"✓ PPTX 生成成功: {output_file}")
        print(f"  总幻灯片数: {len(self.slides)}")
        print(f"  文件大小: {Path(output_file).stat().st_size / 1024:.1f} KB")
        print(f"\n✨ 所有文字完全可编辑，可用 PowerPoint/WPS 打开进行修改")


if __name__ == '__main__':
    converter = MarkdownToPPT()
    input_file = r'F:\hotel-booking-platform\backend-design-defense.md'
    output_file = r'F:\hotel-booking-platform\backend-design-defense.pptx'
    
    converter.generate(input_file, output_file)
